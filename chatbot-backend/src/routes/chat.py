# src/routes/chat.py
from flask import Blueprint, jsonify, request
import google.generativeai as genai
import os
import PyPDF2
from io import BytesIO, StringIO
import uuid
from datetime import datetime
from dotenv import load_dotenv
from pptx import Presentation
import csv
from openpyxl import load_workbook

from src.models.user import db
from src.models.auth import UserSession
from src.models.chat import ChatSession, ChatMessage

load_dotenv()

chat_bp = Blueprint('chat', __name__)

# --- Gemini ---
genai.configure(api_key=os.getenv('GEMINI_API_KEY'))
model = genai.GenerativeModel('gemini-2.5-flash')

# Store document content in memory per-session (ephemeral, in-process only)
# Keyed by session_id so a new chat does not inherit a previous chat's document
session_doc_content = {}

MAX_CONTEXT_CHARS = 10000

def _truncate(text: str, limit: int = MAX_CONTEXT_CHARS) -> str:
    if not text:
        return ""
    return text if len(text) <= limit else text[:limit] + "\n..."

def _extract_pdf_text(file_storage) -> tuple[str, int]:
    reader = PyPDF2.PdfReader(BytesIO(file_storage.read()))
    text_content = ""
    for page in reader.pages:
        page_text = page.extract_text() or ""
        text_content += page_text + "\n"
    return text_content, len(reader.pages)

def _extract_csv_text(file_storage) -> tuple[str, dict]:
    # Read bytes and decode to text for csv reader
    data_bytes = file_storage.read()
    try:
        text = data_bytes.decode('utf-8', errors='replace')
    except Exception:
        text = data_bytes.decode('latin-1', errors='replace')

    f = StringIO(text)
    reader = csv.reader(f)
    rows = []
    header = None
    row_count = 0
    max_preview_rows = 5
    for i, row in enumerate(reader):
        if i == 0:
            header = row
        else:
            if len(rows) < max_preview_rows:
                rows.append(row)
        row_count += 1
        if len("\n".join([",	".join(r) for r in rows])) > MAX_CONTEXT_CHARS:
            break

    columns = len(header) if header else 0
    summary_lines = [
        "CSV Summary:",
        f"Rows (including header): {row_count}",
        f"Columns: {columns}",
        "\nHeader:",
        ", ".join(header or [])
    ]
    if rows:
        summary_lines.append("\nHead (first 5 data rows):")
        for r in rows:
            summary_lines.append(", ".join(r))
    return "\n".join(summary_lines), {"rows": max(row_count - 1, 0), "columns": columns}

def _extract_excel_text(file_storage) -> tuple[str, dict]:
    # Load workbook in read-only mode
    wb = load_workbook(filename=BytesIO(file_storage.read()), read_only=True, data_only=True)
    sheet = wb.active
    sheet_name = sheet.title
    max_rows = sheet.max_row or 0
    max_cols = sheet.max_column or 0

    # Header from first row
    header = []
    if max_rows >= 1:
        header = [str(cell.value) if cell.value is not None else "" for cell in next(sheet.iter_rows(min_row=1, max_row=1, values_only=False))]

    # First 5 data rows
    preview_lines = []
    if max_rows >= 2:
        data_rows_iter = sheet.iter_rows(min_row=2, max_row=min(6, max_rows), values_only=True)
        for row in data_rows_iter:
            preview_lines.append(", ".join([str(v) if v is not None else "" for v in row]))

    summary_lines = [
        f"Excel Summary: Sheet '{sheet_name}'",
        f"Rows (including header): {max_rows}",
        f"Columns: {max_cols}",
        "\nHeader:",
        ", ".join(header)
    ]
    if preview_lines:
        summary_lines.append("\nHead (first 5 data rows):")
        summary_lines.extend(preview_lines)

    text = "\n".join(summary_lines)
    return text, {"sheet": sheet_name, "rows": max(max_rows - 1, 0), "columns": max_cols}

def _extract_ppt_text(file_storage) -> tuple[str, int]:
    prs = Presentation(BytesIO(file_storage.read()))
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"Slide {i}:")
        # Gather text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
        if len("\n".join(lines)) > MAX_CONTEXT_CHARS:
            break
    return "\n".join(lines), len(prs.slides)


@chat_bp.route('/chat', methods=['POST'])
def chat():
    """Chat endpoint: ensures a non-null session_id and persists both messages."""
    try:
        data = request.get_json(silent=True) or {}
        user_message = (data.get('message') or '').strip()
        client_session_id = (data.get('session_id') or '').strip()

        if not user_message:
            return jsonify({'error': 'Message is required'}), 400

        # Always guarantee a session_id
        session_id = client_session_id or str(uuid.uuid4())

        # Fetch or create the chat session
        chat_session = ChatSession.query.filter_by(session_id=session_id).first()
        if not chat_session:
            title = (user_message[:77] + '...') if len(user_message) > 80 else user_message
            # Attach to current user if authenticated
            session_token = request.headers.get('Authorization')
            user_id = None
            if session_token:
                user_session = UserSession.query.filter_by(session_token=session_token, is_active=True).first()
                if user_session:
                    user_id = user_session.user_id

            chat_session = ChatSession(session_id=session_id, user_id=user_id, title=title)
            db.session.add(chat_session)
            db.session.flush()  # ensure session exists before adding messages

        # Determine any document context for this session
        doc_content_for_session = session_doc_content.get(session_id, "")

        # Save the user message
        user_msg = ChatMessage(
            session_id=session_id,
            message_type='user',
            content=user_message,
            has_pdf_context=bool(doc_content_for_session),
        )
        db.session.add(user_msg)

        # Compose prompt/context
        if doc_content_for_session:
            prompt = (
                "Use the following document content to answer.\n\nDOCUMENT:\n"
                f"{doc_content_for_session}\n\nUser question:\n{user_message}"
            )
        else:
            prompt = user_message

        # Call Gemini robustly
        try:
            resp = model.generate_content(prompt)
            bot_text = getattr(resp, 'text', None) or "Sorry, I couldn't generate a response."
            
        except Exception as llm_err:
            bot_text = f"Model error: {llm_err}"

        # Save the bot message
        bot_msg = ChatMessage(
            session_id=session_id,
            message_type='bot',
            content=bot_text,
            has_pdf_context=bool(doc_content_for_session),
        )
        db.session.add(bot_msg)

        # Update session timestamp
        chat_session.updated_at = datetime.utcnow()

        db.session.commit()

        return jsonify({
            'response': bot_text,
            'has_pdf_context': bool(doc_content_for_session),
            'session_id': session_id
        })

    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@chat_bp.route('/chat/sessions', methods=['GET'])
def get_chat_sessions():
    try:
        # If user is authenticated, filter sessions by user_id
        session_token = request.headers.get('Authorization')
        user_id = None
        if session_token:
            user_session = UserSession.query.filter_by(session_token=session_token, is_active=True).first()
            if user_session:
                user_id = user_session.user_id

        query = ChatSession.query
        if user_id:
            query = query.filter_by(user_id=user_id)
        else:
            # For anonymous users, only return sessions with null user_id
            query = query.filter(ChatSession.user_id.is_(None))

        sessions = query.order_by(ChatSession.updated_at.desc()).all()
        return jsonify([s.to_dict() for s in sessions])
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@chat_bp.route('/chat/sessions/<session_id>', methods=['GET'])
def get_chat_session(session_id):
    try:
        session = ChatSession.query.filter_by(session_id=session_id).first()
        if not session:
            return jsonify({'error': 'Session not found'}), 404

        messages = ChatMessage.query.filter_by(session_id=session_id)\
                                    .order_by(ChatMessage.timestamp.asc()).all()

        return jsonify({
            'session': session.to_dict(),
            'messages': [m.to_dict() for m in messages]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@chat_bp.route('/chat/sessions/<session_id>', methods=['DELETE'])
def delete_chat_session(session_id):
    try:
        session_token = request.headers.get('Authorization')
        user_id = None
        if session_token:
            user_session = UserSession.query.filter_by(session_token=session_token, is_active=True).first()
            if user_session:
                user_id = user_session.user_id

        session = ChatSession.query.filter_by(session_id=session_id).first()
        if not session:
            return jsonify({'error': 'Session not found'}), 404

        # Only allow deletion if the session belongs to the current user or is anonymous
        if session.user_id and session.user_id != user_id:
            return jsonify({'error': 'Forbidden'}), 403

        db.session.delete(session)
        db.session.commit()
        return jsonify({'message': 'Session deleted successfully'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@chat_bp.route('/chat/new-session', methods=['POST'])
def create_new_session():
    try:
        session_token = request.headers.get('Authorization')
        user_id = None
        if session_token:
            user_session = UserSession.query.filter_by(session_token=session_token, is_active=True).first()
            if user_session:
                user_id = user_session.user_id

        new_session_id = str(uuid.uuid4())
        chat_session = ChatSession(session_id=new_session_id, user_id=user_id, title='New Chat')
        db.session.add(chat_session)
        db.session.commit()
        return jsonify({'session_id': new_session_id})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@chat_bp.route('/upload-file', methods=['POST'])
def upload_file():
    try:
        # session_id can be sent in form-data or as a query param
        session_id = (request.form.get('session_id')
                      or request.args.get('session_id')
                      or '').strip()

        # If no session_id provided, create one and ensure a session row exists
        if not session_id:
            session_id = str(uuid.uuid4())
            chat_session = ChatSession.query.filter_by(session_id=session_id).first()
            if not chat_session:
                chat_session = ChatSession(session_id=session_id, title='Document Chat')
                db.session.add(chat_session)
                db.session.flush()

        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400

        file = request.files['file']
        if not file.filename:
            return jsonify({'error': 'No file selected'}), 400
        name_lower = file.filename.lower()
        result = { 'session_id': session_id }

        if name_lower.endswith('.pdf'):
            text_content, pages = _extract_pdf_text(file)
            result.update({'message': 'File uploaded successfully', 'kind': 'pdf', 'pages': pages})
        elif name_lower.endswith('.csv'):
            text_content, meta = _extract_csv_text(file)
            result.update({'message': 'File uploaded successfully', 'kind': 'csv', **meta})
        elif name_lower.endswith('.xlsx'):
            text_content, meta = _extract_excel_text(file)
            result.update({'message': 'File uploaded successfully', 'kind': 'xlsx', **meta})
        elif name_lower.endswith('.ppt') or name_lower.endswith('.pptx'):
            text_content, slides = _extract_ppt_text(file)
            result.update({'message': 'File uploaded successfully', 'kind': 'pptx', 'slides': slides})
        else:
            return jsonify({'error': 'Unsupported file type. Allowed: .pdf, .csv, .xlsx, .ppt, .pptx'}), 400

        # Store document text for this session only (truncated)
        text_content = _truncate(text_content)
        session_doc_content[session_id] = text_content

        preview = (text_content[:200] + "...") if len(text_content) > 200 else text_content
        result['preview'] = preview
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@chat_bp.route('/clear-file', methods=['POST'])
def clear_file():
    data = request.get_json(silent=True) or {}
    session_id = (data.get('session_id') or request.args.get('session_id') or '').strip()

    # If a session_id is provided, clear only that session's PDF content
    if session_id:
        session_doc_content.pop(session_id, None)
        return jsonify({'message': 'Document content cleared for session', 'session_id': session_id})

    # Fallback: clear all (avoids stale state, but should not usually be needed)
    session_doc_content.clear()
    return jsonify({'message': 'All document content cleared'})


# Deprecated routes kept for backward compatibility with existing frontends
@chat_bp.route('/upload-pdf', methods=['POST'])
def upload_pdf():
    # Delegate to /upload-file for PDF case only
    return upload_file()


@chat_bp.route('/clear-pdf', methods=['POST'])
def clear_pdf():
    return clear_file()
